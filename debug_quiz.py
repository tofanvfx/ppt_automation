import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

prs = Presentation('template.pptx')

for layout in prs.slide_layouts:
    if 'homework' in layout.name.lower():
        print(f"\n=== Layout: {layout.name} ===")
        print(f"  Shapes: {len(layout.shapes)}")
        for shape in layout.shapes:
            txt = ''
            if hasattr(shape, 'text'):
                txt = repr(shape.text[:80])
            is_ph = getattr(shape, 'is_placeholder', False)
            ph_info = ''
            if is_ph:
                ph_info = f' ph_type={shape.placeholder_format.type} ph_idx={shape.placeholder_format.idx}'
            
            has_img = False
            from lxml import etree
            xml_str = etree.tostring(shape.element, pretty_print=True).decode()
            if 'blip' in xml_str.lower():
                has_img = True
            
            print(f"  - name='{shape.name}' type={shape.shape_type} hasImage={has_img}{ph_info}")
            if txt:
                print(f"    text={txt}")
            
            if shape.shape_type == 6:
                print(f"    Group with {len(shape.shapes)} sub-shapes:")
                for sub in shape.shapes:
                    sub_txt = repr(sub.text[:60]) if hasattr(sub, 'text') else 'N/A'
                    print(f"      - name='{sub.name}' text={sub_txt}")
