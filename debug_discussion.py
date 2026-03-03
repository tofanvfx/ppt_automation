from pptx import Presentation

prs = Presentation('template.pptx')

for layout in prs.slide_layouts:
    if 'discussion' in layout.name.lower():
        print(f"=== Layout: {layout.name} ===")
        for shape in layout.shapes:
            if shape.shape_type == 6: # GROUP
                texts = []
                for sub in shape.shapes:
                    if getattr(sub, 'has_text_frame', False):
                        texts.append(sub.text.replace('\n', ' '))
                print(f"  Group ({len(shape.shapes)} shapes) texts: {texts}")
            else:
                has_text = getattr(shape, 'has_text_frame', False)
                text = shape.text.replace('\n', ' ') if has_text else '<no text>'
                print(f"  {shape.name} (type: {shape.shape_type}, text: {has_text}): {text}")
