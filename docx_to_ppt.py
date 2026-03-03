import re
import os
import copy
import io
from docx import Document
from docx.document import Document as _Document
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
from PIL import Image

ppt_template = "template.pptx"
docx_file = "content.docx"

prs = Presentation(ppt_template)
doc = Document(docx_file)

def clean(text):
    return text.strip().lower()

def get_text(entry):
    """Extract plain text from a content entry (either a string or a (text, ilvl) tuple)."""
    return entry[0] if isinstance(entry, tuple) else entry

def copy_image_rels(element, source_part, target_part):
    """Copy image relationships referenced by blip elements from source to target part."""
    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    r_ns = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'
    for blip in element.findall('.//a:blip', namespaces=ns):
        old_rid = blip.get(f'{r_ns}embed')
        if old_rid and old_rid in source_part.rels:
            rel = source_part.rels[old_rid]
            # Add the image to the target part and get a new rId
            new_rid = target_part.relate_to(rel.target_part, rel.reltype)
            blip.set(f'{r_ns}embed', new_rid)

def get_layout(name):
    for layout in prs.slide_layouts:
        if layout.name == f"LAYOUT_{name}":
            return layout
            
    # Try case-insensitive, underscore-insensitive match
    target = f"layout_{name}".replace("_", "").lower()
    for layout in prs.slide_layouts:
        if layout.name.replace("_", "").replace(" ", "").lower() == target:
            return layout
            
    if name.replace("_", "").lower() == 'mathtitlepage':
        for layout in prs.slide_layouts:
            if layout.name == 'LAYOUT_math_page_title':
                return layout

    if name.replace("_", "").lower() == 'ssttitlepage':
        for layout in prs.slide_layouts:
            if layout.name == 'LAYOUT_sst_page_title':
                return layout
                
    if name.replace("_", "").lower() in ('learningobjective', 'sstlopage', '1layoutsstlopage'):
        for layout in prs.slide_layouts:
            if layout.name == '1_LAYOUT_sst_lo_page' or layout.name == 'LAYOUT_sst_lo_page':
                return layout

    if name.replace("_", "").lower() == 'mathlopage':
        for layout in prs.slide_layouts:
            if layout.name == 'LAYOUT_math_lo_page':
                return layout

    if name.replace("_", "").lower() in ('finalquizpage', 'finalquiz', 'quiz'):
        for layout in prs.slide_layouts:
            if layout.name == 'LAYOUT_final_quiz_page':
                return layout

    if name.replace("_", "").lower() in ('quiztimepage', 'sstquiztimepage', 'sst_quiztime_page'):
        # Return _01 as default; the main loop will override based on option lengths
        for layout in prs.slide_layouts:
            if layout.name == 'LAYOUT_sst_quiztime_page_01':
                return layout

    if name.replace("_", "").lower() in ('sstcontentpage01', 'sstcontentpage1'):
        for layout in prs.slide_layouts:
            if layout.name == 'LAYOUT_sst_content_page_01':
                return layout

    if name.replace("_", "").lower() in ('sstsummarypage', 'sstsummary', 'sstsummerypage', 'sstsummery'):
        for layout in prs.slide_layouts:
            if layout.name == 'LAYOUT_sst_summary_page':
                return layout

    if name.replace("_", "").lower() in ('mathsummarypage', 'mathsummary', 'mathsummerypage', 'mathsummery'):
        for layout in prs.slide_layouts:
            if layout.name == 'LAYOUT_math_summary_page':
                return layout

    if name.replace("_", "").lower() in ('notedownpage', 'sstnotedownpage', 'sst_notedown_page'):
        for layout in prs.slide_layouts:
            if layout.name == 'LAYOUT_sst_notedown_page':
                return layout

    if name.replace("_", "").lower() in ('previouspage', 'sstpreviouspage', 'sst_previous_page'):
        for layout in prs.slide_layouts:
            if layout.name == 'LAYOUT_sst_previous_page':
                return layout

    if name.replace("_", "").lower() in ('homeworkpage', 'ssthomeworkpage', 'ssthomework', 'homework'):
        for layout in prs.slide_layouts:
            if layout.name == 'LAYOUT_sst_homework_page':
                return layout

    if name.replace("_", "").lower() in ('discussionpage', 'sstdiscussionpage', 'sstdiscussion', 'sst_discussion_page'):
        for layout in prs.slide_layouts:
            if layout.name == 'LAYOUT_sst_discussion_page':
                return layout

    return None

# Nested dictionary to store templates per layout
sst_content_templates = {
    'LAYOUT_sst_content_page_01': {'topic': None, 'subtopic': None, 'text': None},
    'LAYOUT_sst_content_page_02': {'topic': None, 'subtopic': None, 'text': None},
    'LAYOUT_sst_notedown_page': {'text': None},
    'LAYOUT_sst_quiztime_page_01': {'title': None, 'question': None, 'options': [], 'picture': None},
    'LAYOUT_sst_quiztime_page_02': {'title': None, 'question': None, 'options': [], 'picture': None},
    'LAYOUT_sst_discussion_page': {'question1': None}
}

# Extract lo_page group shape XML before processing slides
lo_group_xmls = {
    '1_LAYOUT_sst_lo_page': None,
    'LAYOUT_sst_lo_page': None,
    'LAYOUT_math_lo_page': None,
    'LAYOUT_sst_summary_page': None,
    'LAYOUT_math_summary_page': None,
    'LAYOUT_sst_previous_page': None
}

# Store static title templates for lo/summary layouts where add_slide fails
lo_title_xmls = {k: None for k in lo_group_xmls.keys()}
lo_title_xmls['1_LAYOUT_sst_lo_page'] = None # Ensure renamed one is included
lo_title_xmls['LAYOUT_sst_previous_page'] = None

for layout in prs.slide_layouts:
    if layout.name in lo_group_xmls:
        for shape in layout.shapes:
            if shape.shape_type == 6:  # msoGroup
                for subshape in shape.shapes:
                    if getattr(subshape, 'has_text_frame', False) and 'Text goes here' in subshape.text:
                        lo_group_xmls[layout.name] = copy.deepcopy(shape.element)
                        shape.element.getparent().remove(shape.element)
                        break
                if lo_group_xmls[layout.name] is not None:
                    break
    
    if layout.name in lo_title_xmls:
        for shape in layout.shapes:
            if shape.shape_type != 6 and getattr(shape, 'has_text_frame', False) and shape.text.strip():
                lo_title_xmls[layout.name] = copy.deepcopy(shape.element)
                shape.element.getparent().remove(shape.element)
                # print(f"DEBUG: Captured title for {layout.name}: '{shape.text[:20]}...'")
                break
    
    if layout.name in sst_content_templates:
        templates = sst_content_templates[layout.name]
        for shape in list(layout.shapes): # Use list() to avoid issues when removing
            has_topic = False
            has_subtopic = False
            has_quiz_title = False
            
            # 0. Capture standalone PICTURE shapes for quiztime layouts
            if layout.name.startswith('LAYOUT_sst_quiztime_page') and shape.shape_type == 13:  # PICTURE
                templates['picture'] = copy.deepcopy(shape.element)
                shape.element.getparent().remove(shape.element)
                continue

            # 1. Check groups for topic/subtopic or "QUIZ TIME"
            if shape.shape_type == 6:  # GROUP
                for sub in shape.shapes:
                    if getattr(sub, 'has_text_frame', False):
                        txt = sub.text.strip().lower()
                        if txt == 'topic': has_topic = True
                        elif txt == 'subtopic': has_subtopic = True
                        elif layout.name.startswith('LAYOUT_sst_quiztime_page') and 'quiz time' in txt:
                            has_quiz_title = True
                
                if has_quiz_title:
                    templates['title'] = copy.deepcopy(shape.element)
                    shape.element.getparent().remove(shape.element)
                    continue # Move to next shape

            # 2. Check individual shapes for topic/subtopic
            elif getattr(shape, 'has_text_frame', False):
                txt = shape.text.strip().lower()
                if txt == 'topic': has_topic = True
                elif txt == 'subtopic': has_subtopic = True
            
            if has_topic:
                templates['topic'] = copy.deepcopy(shape.element)
                shape.element.getparent().remove(shape.element)
            elif has_subtopic:
                templates['subtopic'] = copy.deepcopy(shape.element)
                shape.element.getparent().remove(shape.element)
            
            # 3. Check for body/placeholder text (question, options, or general text)
            elif shape.has_text_frame and ('Text goes here' in shape.text or (shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.BODY)):
                txt = shape.text.lower()
                if layout.name.startswith('LAYOUT_sst_quiztime_page'):
                    if 'quiz time' in txt:
                        templates['title'] = copy.deepcopy(shape.element)
                    elif 'question' in txt:
                        templates['question'] = copy.deepcopy(shape.element)
                    elif 'options' in txt:
                        templates['options'].append(copy.deepcopy(shape.element))
                    elif txt.strip(): 
                        if not templates['title']: templates['title'] = copy.deepcopy(shape.element)
                elif layout.name == 'LAYOUT_sst_discussion_page':
                    if 'question1' in txt:
                        templates['question1'] = copy.deepcopy(shape.element)
                else:
                    templates['text'] = copy.deepcopy(shape.element)
                
                if shape.element.getparent() is not None:
                    shape.element.getparent().remove(shape.element)

def iter_block_items(parent):
    if isinstance(parent, _Document):
        parent_elm = parent.element.body
    elif isinstance(parent, _Cell):
        parent_elm = parent._tc
    else:
        raise ValueError("iter_block_items: parent must be Document or _Cell")

    for child in parent_elm.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, parent)
        elif isinstance(child, CT_Tbl):
            yield Table(child, parent)


# Parse DOCX into sections
sections = []
current_section = None

w_ns = 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'

for item in iter_block_items(doc):
    lines_to_process = []
    images_to_add = []
    
    if isinstance(item, Paragraph):
        text = item.text.strip()
        # Detect indent level from numPr/ilvl
        ilvl = 0
        numPr = item._element.find(f'.//{{{w_ns}}}numPr')
        if numPr is not None:
            ilvl_elem = numPr.find(f'{{{w_ns}}}ilvl')
            if ilvl_elem is not None:
                ilvl = int(ilvl_elem.get(f'{{{w_ns}}}val', '0'))
        lines_to_process.append((text, ilvl))
        # Extract images from paragraph
        for run in item.runs:
            for drawing in run._element.findall('.//w:drawing', namespaces=run._element.nsmap):
                for blip in drawing.findall('.//a:blip', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}):
                    embed = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                    if embed:
                        blob = doc.part.related_parts[embed].blob
                        images_to_add.append(blob)
    elif isinstance(item, Table):
        for row in item.rows:
            # Simple conversion: join first two cells if they exist with ":"
            cells = [c.text.strip() for c in row.cells]
            if len(cells) >= 2:
                key = cells[0].rstrip(':').strip()
                val = cells[1].lstrip(':').strip()
                lines_to_process.append((f"{key}: {val}", 0))
            elif len(cells) == 1:
                lines_to_process.append((cells[0], 0))

    for entry in lines_to_process:
        text = entry[0] if isinstance(entry, tuple) else entry
        ilvl = entry[1] if isinstance(entry, tuple) else 0
        
        if not text and not images_to_add:
            continue
            
        match = re.search(r'\[\s*([^\]]+?)\s*\]', text)
        if match:
            name = match.group(1).strip()
            current_section = {
                'name': name,
                'content': [],
                'images': []
            }
            sections.append(current_section)
        elif current_section is not None:
            if text:
                current_section['content'].append((text, ilvl))
            if images_to_add:
                current_section['images'].extend(images_to_add)
                images_to_add = []

def replace_text_preserve_format(shape, new_text, center=False, font_color=None):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    if not tf.paragraphs or not tf.paragraphs[0].runs:
        if isinstance(new_text, list):
            shape.text = "\n".join(new_text)
        else:
            shape.text = str(new_text)
        if center:
            for p in tf.paragraphs: p.alignment = PP_ALIGN.CENTER
        return
        
    p0 = tf.paragraphs[0]
    p0_xml = copy.deepcopy(p0._p)
    
    font = p0.runs[0].font
    name = font.name
    size = font.size
    bold = font.bold
    italic = font.italic
    underline = font.underline

    try:
        color_type = getattr(font.color, 'type', None)
        color_rgb = getattr(font.color, 'rgb', None) if color_type == 1 else None
        color_theme = getattr(font.color, 'theme_color', None) if color_type == 2 else None
    except:
        color_type, color_rgb, color_theme = None, None, None

    tf.clear()
    
    texts = new_text if isinstance(new_text, list) else [str(new_text)]
    
    for i, text in enumerate(texts):
        if i == 0:
            p = tf.paragraphs[0]
            new_p_xml = copy.deepcopy(p0_xml)
            p._p.getparent().replace(p._p, new_p_xml)
            p = tf.paragraphs[0]
        else:
            new_p_xml = copy.deepcopy(p0_xml)
            tf._txBody.append(new_p_xml)
            p = tf.paragraphs[i]
            p.space_before = Pt(12)
            
        # Aggressively clear existing runs and fields in the paragraph XML
        for r_elem in p._p.findall('.//a:r', namespaces=p._p.nsmap):
            p._p.remove(r_elem)
        for fld_elem in p._p.findall('.//a:fld', namespaces=p._p.nsmap):
            p._p.remove(fld_elem)
            
        if center:
            p.alignment = PP_ALIGN.CENTER
            
        new_run = p.add_run()
        new_run.text = str(text)
        
        if name is not None: new_run.font.name = name
        if size is not None: new_run.font.size = size
        if bold is not None: new_run.font.bold = bold
        if italic is not None: new_run.font.italic = italic
        if underline is not None: new_run.font.underline = underline
        
        if font_color is not None:
            new_run.font.color.rgb = font_color
        else:
            if color_rgb is not None:
                new_run.font.color.rgb = color_rgb
            elif color_theme is not None:
                new_run.font.color.theme_color = color_theme

for section in sections:
    sname = section['name'].strip().lower()
    if sname == 'sst_content_page':
        if len(section.get('images', [])) > 1:
            layout = get_layout('sst_content_page_02')
        else:
            layout = get_layout('sst_content_page_01')
    elif sname in ('sst_quiztime_page', 'quiztime_page', 'sstquiztimepage'):
        # Parse quiz content first to determine layout
        quiz_data = {'question': '', 'options': []}
        for entry in section['content']:
            line = get_text(entry)
            line_low = line.lower()
            if 'question:' in line_low:
                quiz_data['question'] = line.split(':', 1)[1].strip()
            elif 'options:' in line_low:
                opts = line.split(':', 1)[1].strip()
                quiz_data['options'] = [o.strip() for o in opts.split(',')]
            elif line.strip():
                if not quiz_data['question']: quiz_data['question'] = line.strip()
                else: quiz_data['options'].append(line.strip())
        
        # Pick layout based on option character length
        # If ANY option > 25 chars → _01 (long), else _02 (short)
        use_long = any(len(opt) > 25 for opt in quiz_data['options'])
        if use_long:
            layout = get_layout('sst_quiztime_page_01')
            print(f"Quiz options have long text (>25 chars) -> using LAYOUT_sst_quiztime_page_01")
        else:
            layout = get_layout('sst_quiztime_page_02')
            print(f"Quiz options are short (<=25 chars) -> using LAYOUT_sst_quiztime_page_02")
    else:
        layout = get_layout(section['name'])

    if not layout:
        print(f"Skipping section [{section['name']}], layout not found.")
        continue
        
    slide = prs.slides.add_slide(layout)
    
    if section['name'].replace("_", "").lower() in ('mathpagetitle', 'mathtitlepage', 'sstpagetitle', 'ssttitlepage'):
        data = {}
        for entry in section['content']:
            line = get_text(entry)
            if ":" in line:
                parts = line.split(":", 1)
                data[parts[0].strip().upper()] = parts[1].strip()
        
        idx_mapping = {}
        for shape in layout.shapes:
            if shape.is_placeholder and shape.has_text_frame:
                idx_mapping[clean(shape.text)] = shape.placeholder_format.idx
                
        mapping = {
            "CLASS": "class",
            "SUBJECT": "subject",
            "CHAPTER_NUMBER": "chapter number",
            "CHAPTER_NAME": "chapter name",
            "LESSON": "lesson",
            "TOPIC": "topic"
        }

        for key, template_word in mapping.items():
            if key in data:
                cleaned_word = clean(template_word)
                if cleaned_word in idx_mapping:
                    idx = idx_mapping[cleaned_word]
                    found_on_slide = False
                    for shape in slide.shapes:
                        if shape.is_placeholder and shape.placeholder_format.idx == idx:
                            shape.text = data[key]
                            print(f"Updated PageTitle: {template_word} -> {data[key]}")
                            found_on_slide = True
                            break
                    if not found_on_slide:
                        print(f"[{template_word}] idx {idx} NOT FOUND on slide shapes!")
                else:
                    print(f"[{template_word}] NOT FOUND in idx_mapping! keys={idx_mapping.keys()}")
            else:
                print(f"[{key}] NOT FOUND in docx data!")
    elif layout.name in lo_group_xmls:
        # Inject static title if available
        if layout.name in lo_title_xmls and lo_title_xmls[layout.name] is not None:
            slide.shapes._spTree.append(copy.deepcopy(lo_title_xmls[layout.name]))

        group_xml = lo_group_xmls.get(layout.name)
            
        if group_xml is not None:
            # Group content: merge sub-bullets (ilvl > 0) with their parent line
            # Each item: (main_text, [(sub_text, ilvl), ...])
            grouped_items = []
            for entry in section['content']:
                text = entry[0] if isinstance(entry, tuple) else entry
                ilvl = entry[1] if isinstance(entry, tuple) else 0
                clean_line = text.strip()
                if clean_line.startswith(('•', '-', '*')):
                    clean_line = clean_line.lstrip('•-*').strip()
                
                if ilvl == 0:
                    grouped_items.append((clean_line, []))
                elif ilvl > 0 and grouped_items:
                    grouped_items[-1][1].append((clean_line, ilvl))
                else:
                    grouped_items.append((clean_line, []))
            
            current_top_offset = 0
            for main_text, sub_entries in grouped_items:
                new_element = copy.deepcopy(group_xml)
                slide.shapes._spTree.append(new_element)
                new_shape = slide.shapes[-1]
                new_shape.top = new_shape.top + current_top_offset
                
                for subshape in new_shape.shapes:
                    if getattr(subshape, 'has_text_frame', False) and 'Text goes here' in subshape.text:
                        tf = subshape.text_frame
                        # Set main text in first paragraph
                        tf.paragraphs[0].text = main_text
                        for run in tf.paragraphs[0].runs:
                            run.font.size = Pt(36)
                            run.font.color.rgb = RGBColor(255, 255, 255)
                        
                        # Add sub-bullet paragraphs with nested indentation
                        for bullet_text, bullet_ilvl in sub_entries:
                            p = tf.add_paragraph()
                            # Indent and bullet marker based on nesting level
                            indent = '    ' * bullet_ilvl
                            markers = {1: '•', 2: '◦', 3: '▪'}
                            marker = markers.get(bullet_ilvl, '▪')
                            p.text = f'{indent}{marker} {bullet_text}'
                            p.space_before = Pt(4)
                            # Smaller font for deeper levels
                            font_size = max(20, 36 - (bullet_ilvl * 4))
                            for run in p.runs:
                                run.font.size = Pt(font_size)
                                run.font.color.rgb = RGBColor(255, 255, 255)
                
                # Calculate height based on content
                chars_per_line = 65
                line_count = max(1, len(main_text) // chars_per_line + (1 if len(main_text) % chars_per_line > 0 else 0))
                sub_height = len(sub_entries) * 0.45
                current_top_offset += int(Inches(0.9 + 0.55 * (line_count - 1) + sub_height))
                
            total_subs = sum(len(s) for _, s in grouped_items)
            print(f"Updated {layout.name} with {len(grouped_items)} items ({total_subs} sub-bullets)")
    elif layout.name in ('LAYOUT_sst_quiztime_page_01', 'LAYOUT_sst_quiztime_page_02'):
        templates = sst_content_templates.get(layout.name, {})
        
        # Inject static picture
        if templates.get('picture') is not None:
            pic_elem = copy.deepcopy(templates['picture'])
            copy_image_rels(pic_elem, layout.part, slide.part)
            slide.shapes._spTree.append(pic_elem)

        # Inject static title
        if templates.get('title') is not None:
            title_elem = copy.deepcopy(templates['title'])
            copy_image_rels(title_elem, layout.part, slide.part)
            slide.shapes._spTree.append(title_elem)

        # Reuse quiz_data if already parsed during layout selection, otherwise parse now
        if 'quiz_data' not in dir():
            quiz_data = {'question': '', 'options': []}
            for entry in section['content']:
                line = get_text(entry)
                line_low = line.lower()
                if 'question:' in line_low:
                    quiz_data['question'] = line.split(':', 1)[1].strip()
                elif 'options:' in line_low:
                    opts = line.split(':', 1)[1].strip()
                    quiz_data['options'] = [o.strip() for o in opts.split(',')]
                elif line.strip():
                    if not quiz_data['question']: quiz_data['question'] = line.strip()
                    else: quiz_data['options'].append(line.strip())

        if templates.get('question') is not None:
            slide.shapes._spTree.append(copy.deepcopy(templates['question']))
            replace_text_preserve_format(slide.shapes[-1], quiz_data['question'])
        
        # Inject options
        option_templates = templates.get('options', [])
        if len(option_templates) == 1:
            # _01 style: single placeholder, all options as multi-line text
            slide.shapes._spTree.append(copy.deepcopy(option_templates[0]))
            replace_text_preserve_format(slide.shapes[-1], quiz_data['options'])
        elif len(option_templates) > 1:
            # _02 style: separate placeholder per option
            for idx, opt_template in enumerate(option_templates):
                slide.shapes._spTree.append(copy.deepcopy(opt_template))
                if idx < len(quiz_data['options']):
                    replace_text_preserve_format(slide.shapes[-1], quiz_data['options'][idx])
                else:
                    replace_text_preserve_format(slide.shapes[-1], '')
        
        print(f"Updated {layout.name} with quiz question and {len(quiz_data['options'])} options.")
        quiz_data = None  # Reset for next quiz section
    elif layout.name in ('LAYOUT_sst_content_page_01', 'LAYOUT_sst_content_page_02'):
        import io
        data = {}
        data_text_list = []
        for entry in section['content']:
            line = get_text(entry)
            if ":" in line:
                parts = line.split(":", 1)
                key = parts[0].strip().lower()
                val = parts[1].strip()
                if key == 'text':
                    data_text_list.append(val)
                else:
                    data[key] = val
            elif line.strip() and not line.lower().startswith(('topic', 'subtopic')):
                data_text_list.append(line.strip())
                
        if data_text_list:
            data['text'] = data_text_list
        
        # Inject extracted XML templates onto this slide from the correct layout version
        shape_elements = []
        templates = sst_content_templates.get(layout.name, {})
        
        if templates.get('topic') is not None:
            topic_elem = copy.deepcopy(templates['topic'])
            slide.shapes._spTree.append(topic_elem)
            shape_elements.append(slide.shapes[-1])
            
        if templates.get('subtopic') is not None and 'subtopic' in data:
            subtopic_elem = copy.deepcopy(templates['subtopic'])
            slide.shapes._spTree.append(subtopic_elem)
            shape_elements.append(slide.shapes[-1])
            
        if templates.get('text') is not None:
            text_elem = copy.deepcopy(templates['text'])
            # If no subtopic, move the text box to where the subtopic would have been
            if 'subtopic' not in data and templates.get('subtopic') is not None:
                sub_xml = templates['subtopic']
                # Search for offset element in subtopic template
                # The namespace for 'a' is usually needed
                ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                sub_off = sub_xml.find('.//a:off', namespaces=ns)
                text_off = text_elem.find('.//a:off', namespaces=ns)
                
                if sub_off is not None and text_off is not None:
                    # Move text to subtopic's Y coordinate
                    text_off.set('y', sub_off.get('y'))
            
            slide.shapes._spTree.append(text_elem)
            shape_elements.append(slide.shapes[-1])
        
        # Insert text and subtexts
        for shape in shape_elements:
            if shape.has_text_frame and 'Text goes here' in shape.text:
                if 'text' in data:
                    replace_text_preserve_format(shape, data['text'], font_color=RGBColor(255, 255, 255))
            
            # Check for topic/subtopic either as standalone shape or inside group
            is_group = (getattr(shape, 'shape_type', None) == 6)
            targets = shape.shapes if is_group else [shape]
            
            for target in targets:
                if getattr(target, 'has_text_frame', False):
                    txt = target.text.strip().lower()
                    if 'topic' in data or 'subtopic' in data:
                        # print(f"DEBUG: Comparing txt='{txt}' with data keys {list(data.keys())}")
                        pass
                    if txt == 'topic' and 'topic' in data:
                        replace_text_preserve_format(target, data['topic'], center=True)
                    elif txt == 'subtopic' and 'subtopic' in data:
                        replace_text_preserve_format(target, data['subtopic'], center=True)
                        target.text_frame.word_wrap = False
                        # Resize logic: if it's a group, resize all children. If standalone, resize it.
                        estimated_width = len(data['subtopic']) * 300000 + 400000
                        if estimated_width > shape.width:
                            diff = estimated_width - shape.width
                            shape.width = estimated_width
                            if is_group:
                                for child in shape.shapes:
                                    child.width = child.width + diff
                            else:
                                pass # Shape width already updated
                            
        # Handle picture placeholder(s)
        images = section.get('images', [])
        pic_ph = None
        for shape in slide.shapes:
            if getattr(shape, 'is_placeholder', False) and shape.placeholder_format.type == 18:
                pic_ph = shape
                break
        
        if pic_ph is not None and images:
            # Capture properties of the placeholder before it's potentially replaced/removed
            left, top, width, height = pic_ph.left, pic_ph.top, pic_ph.width, pic_ph.height
            num_images = len(images)
            
            if num_images > 1:
                # Remove the placeholder to avoid overlap with newly added pictures
                pic_ph.element.getparent().remove(pic_ph.element)
                
                gap = Inches(0.1)
                total_gap = gap * (num_images - 1)
                img_width = (width - total_gap) / num_images
                
                for i, img_blob in enumerate(images):
                    # Load image to get original dimensions for aspect ratio
                    with Image.open(io.BytesIO(img_blob)) as pil_img:
                        orig_w, orig_h = pil_img.size
                    
                    slot_w = img_width
                    slot_h = height
                    
                    # Calculate fit dimensions
                    slot_aspect = slot_w / slot_h
                    img_aspect = orig_w / orig_h
                    
                    if img_aspect > slot_aspect:
                        # Image is wider than slot relative to height
                        draw_w = slot_w
                        draw_h = slot_w / img_aspect
                    else:
                        # Image is taller than slot relative to width
                        draw_h = slot_h
                        draw_w = slot_h * img_aspect
                    
                    # Center within slot
                    final_left = int(left + (i * (slot_w + gap)) + (slot_w - draw_w) / 2)
                    final_top = int(top + (slot_h - draw_h) / 2)
                    
                    slide.shapes.add_picture(
                        io.BytesIO(img_blob), 
                        final_left, 
                        final_top, 
                        width=int(draw_w),
                        height=int(draw_h)
                    )
            else:
                # Single image: use standard placeholder insertion (handles aspect ratio better)
                pic_ph.insert_picture(io.BytesIO(images[0]))
        elif pic_ph is not None:
            # Hide/Remove placeholder if no images are present
            pic_ph.element.getparent().remove(pic_ph.element)
                    
        print(f"Updated {layout.name} with text and {len(images)} images.")

    elif layout.name == 'LAYOUT_final_quiz_page':
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                shape.text = "Final Quiz!"
                break
        print("Inserted LAYOUT_final_quiz_page (Welcome Slide)")

        qa_pairs = []
        for entry in section['content']:
            line = get_text(entry)
            if ":" in line:
                parts = line.split(":", 1)
                q_text = parts[0].strip()
                a_text = parts[1].strip()
                if q_text.lower() == 'question' and a_text.lower() == 'answer':
                    continue
                qa_pairs.append((q_text, a_text))
        
        layout_q = get_layout('final_quiz_page_q')
        layout_a = get_layout('final_quiz_page_a')
        
        for q, a in qa_pairs:
            if layout_q:
                s_q = prs.slides.add_slide(layout_q)
                for shape in s_q.shapes:
                    if shape.is_placeholder and shape.placeholder_format.idx == 0:
                        shape.text = q
            if layout_a:
                s_a = prs.slides.add_slide(layout_a)
                for shape in s_a.shapes:
                    if shape.is_placeholder and shape.placeholder_format.idx == 0:
                        shape.text = a
        print(f"Generated {len(qa_pairs)} Question and Answer slide pairs.")
    elif layout.name == 'LAYOUT_sst_notedown_page':
        templates = sst_content_templates.get(layout.name, {})
        if templates.get('text') is not None:
            text_elem = copy.deepcopy(templates['text'])
            slide.shapes._spTree.append(text_elem)
            shape = slide.shapes[-1]
            
            processed_content = []
            for entry in section['content']:
                line = get_text(entry)
                clean_line = line.strip()
                if clean_line.startswith(("•", "-", "*")):
                    clean_line = "❖ " + clean_line.lstrip("•-*").strip()
                processed_content.append(clean_line)
            
            replace_text_preserve_format(shape, processed_content)
            print(f"Updated {layout.name} Body text with preserved formatting and spacing.")

    elif layout.name == 'LAYOUT_sst_discussion_page':
        templates = sst_content_templates.get(layout.name, {})
        if templates.get('question1') is not None:
            # Parse 'question1: <text>' from the content section
            q1_text = ""
            for entry in section['content']:
                line = get_text(entry)
                if ":" in line:
                    parts = line.split(":", 1)
                    if parts[0].strip().lower() == 'question1':
                        q1_text = parts[1].strip()
                        break
                elif line.strip() and not q1_text:
                    # Fallback: if no colon, just use the first non-empty line
                    q1_text = line.strip()
            
            q1_elem = copy.deepcopy(templates['question1'])
            slide.shapes._spTree.append(q1_elem)
            shape = slide.shapes[-1]
            if q1_text:
                replace_text_preserve_format(shape, q1_text)
                print(f"Updated {layout.name} question1 -> '{q1_text[:20]}...'")
            else:
                replace_text_preserve_format(shape, "question1")  # Keep original if missing
                print(f"Inserted {layout.name} with default text")

if __name__ == "__main__":
    import sys
    if hasattr(sys.stdout, 'reconfigure'):
        sys.stdout.reconfigure(encoding='utf-8')

    try:
        prs.save("Generated_Presentation.pptx")
        print("DONE -- PPT Generated (Generated_Presentation.pptx)")
    except PermissionError:
        prs.save("Generated_Presentation_new.pptx")
        print("DONE -- PPT Generated (Saved as Generated_Presentation_new.pptx because original is open)")