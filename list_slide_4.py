from pptx import Presentation

prs = Presentation("Generated_Presentation.pptx")
slide = prs.slides[3] # Slide 4

print(f"Slide 4 Layout: {slide.slide_layout.name}")
for i, shape in enumerate(slide.shapes):
    print(f"Shape {i}: '{shape.name}' | Type: {shape.shape_type}")
    if shape.has_text_frame:
        print(f"  - Text: '{shape.text}'")
